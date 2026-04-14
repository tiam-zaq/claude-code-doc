"""
PPTXファイルから指定したスライドの内容を抽出するスクリプト。
カンマ区切りで複数スライドを指定可能。

Usage:
    python extract_slide.py <pptx_path> <slide_spec>

    slide_spec:
        "2"       → スライド2を抽出
        "2,4,8"   → スライド2, 4, 8を抽出
        "概要"     → タイトルに「概要」を含むスライドを検索
"""

import sys
from pathlib import Path


def extract_one_slide(prs, slide_spec: str) -> dict:
    total_slides = len(prs.slides)
    spec = slide_spec.strip()

    target_slide = None
    target_index = None

    if spec.isdigit():
        idx = int(spec) - 1
        if 0 <= idx < total_slides:
            target_slide = prs.slides[idx]
            target_index = idx + 1
        else:
            print(f"WARNING: スライド番号 {spec} は範囲外です（全{total_slides}枚）→ スキップ", file=sys.stderr)
            return None
    else:
        for i, slide in enumerate(prs.slides):
            if slide.shapes.title and spec in slide.shapes.title.text:
                target_slide = slide
                target_index = i + 1
                break
        if target_slide is None:
            print(f"WARNING: タイトルに '{spec}' を含むスライドが見つかりません → スキップ", file=sys.stderr)
            return None

    title_text = ""
    if target_slide.shapes.title and target_slide.shapes.title.text.strip():
        title_text = target_slide.shapes.title.text.strip()

    body_texts = []
    for shape in target_slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text.strip()
        if not text or text == title_text:
            continue
        lines = [line.strip() for line in text.splitlines() if line.strip()]
        if lines:
            body_texts.append("\n".join(lines))

    return {
        "file": Path(prs.part.partname).name if hasattr(prs.part, 'partname') else "",
        "slide_number": target_index,
        "total_slides": total_slides,
        "title": title_text or f"スライド {target_index}",
        "body": "\n\n".join(body_texts),
    }


def format_multi_slides(slides: list, pptx_name: str, total: int) -> str:
    lines = [
        "## 参照スライド情報",
        "",
        f"- **ファイル**: {pptx_name}",
        f"- **参照ページ**: {', '.join(str(s['slide_number']) for s in slides)} / 全{total}枚",
        "",
    ]

    for s in slides:
        lines.extend([
            f"---",
            f"",
            f"### P{s['slide_number']}: {s['title']}",
            f"",
            s["body"] if s["body"] else "（本文テキストなし）",
            f"",
        ])

    lines.extend([
        "---",
        "",
        "> このIssueはスライド内容から自動生成されました。",
        "> 上記テーマについてAIエージェントが多角的に議論します。",
    ])
    return "\n".join(lines)


if __name__ == "__main__":
    if len(sys.argv) != 3:
        print("Usage: python extract_slide.py <pptx_path> <slide_spec>")
        sys.exit(1)

    try:
        from pptx import Presentation
    except ImportError:
        print("ERROR: python-pptx が必要です。pip install python-pptx を実行してください。")
        sys.exit(1)

    pptx_path = sys.argv[1]
    slide_spec = sys.argv[2]

    path = Path(pptx_path)
    if not path.exists():
        print(f"ERROR: ファイルが見つかりません: {pptx_path}")
        sys.exit(1)

    prs = Presentation(str(path))

    # カンマ区切りで複数スライド対応
    specs = [s.strip() for s in slide_spec.split(",") if s.strip()]
    slides = []
    for spec in specs:
        result = extract_one_slide(prs, spec)
        if result:
            slides.append(result)

    if not slides:
        print("ERROR: 有効なスライドが1枚も見つかりませんでした")
        sys.exit(1)

    total = len(prs.slides)
    print(format_multi_slides(slides, path.name, total))

    first_title = slides[0]["title"]
    if len(slides) > 1:
        first_title += f" 他{len(slides)-1}枚"
    print(f"\n__SLIDE_TITLE__={first_title}", file=sys.stderr)
